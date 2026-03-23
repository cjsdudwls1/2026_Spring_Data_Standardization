"""EDA 보고서를 발표용 PPT로 변환하는 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 설정 ──
FIGURES = "c:/cheon/cheon_wokespace/2주차 과제/output/figures"
OUT = "c:/cheon/cheon_wokespace/2주차 과제/EDA_발표자료_v2.pptx"

# 색상 팔레트
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # 짙은 남색 배경
ACCENT  = RGBColor(0x00, 0xB4, 0xD8)        # 청록 포인트
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xE0, 0xE0, 0xE0)
YELLOW  = RGBColor(0xFF, 0xD6, 0x00)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
GREEN   = RGBColor(0x00, 0xC9, 0xA7)
RED     = RGBColor(0xFF, 0x4D, 0x6D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── 유틸리티 함수 ──
def add_bg(slide, color=DARK_BG):
    """슬라이드 배경 채우기"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, size=16, color=LIGHT, bold=False, space_before=Pt(4),
             font_name="맑은 고딕", alignment=PP_ALIGN.LEFT):
    """텍스트 프레임에 문단 추가"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.alignment = alignment
    return p


def add_image_safe(slide, path, left, top, width=None, height=None):
    """이미지 안전 추가"""
    full = os.path.join(FIGURES, path)
    if os.path.exists(full):
        kwargs = {}
        if width:  kwargs['width']  = Inches(width)
        if height: kwargs['height'] = Inches(height)
        slide.shapes.add_picture(full, Inches(left), Inches(top), **kwargs)
        return True
    return False


def add_accent_bar(slide, left=0.4, top=1.05, width=1.8, height=0.06):
    """포인트 컬러 바"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_section_number(slide, number, left=0.4, top=0.25):
    """섹션 번호 원형 배지"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                    Inches(left), Inches(top),
                                    Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def make_slide():
    """빈 슬라이드 생성 + 배경"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    return slide


def add_table(slide, rows, cols, data, left, top, width, height,
              header_color=ACCENT, cell_color=LIGHT, font_size=11):
    """테이블 추가"""
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "맑은 고딕"
                if i == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = RGBColor(0x1B, 0x1B, 0x2F)
            # 헤더 행 배경
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF5)
    return table


# ════════════════════════════════════════
# 슬라이드 1: 표지
# ════════════════════════════════════════
slide = make_slide()
# 상단 장식 바
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0),
                                SLIDE_W, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, 1.5, 1.8, 10, 1.2,
         "혈액 기부 서비스 센터 데이터", 42, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, 1.5, 2.8, 10, 1.0,
         "탐색적 데이터 분석 (EDA) 보고서", 36, WHITE, True, PP_ALIGN.CENTER)

# 구분선
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(4.5), Inches(4.0),
                                Inches(4.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, 1.5, 4.3, 10, 0.5,
         "Statistics for Machine Learning", 20, LIGHT, False, PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.9, 10, 0.5,
         "데이터셋: OpenML - blood-transfusion-service-center (ID: 1464)", 16, LIGHT, False, PP_ALIGN.CENTER)
add_text(slide, 1.5, 5.4, 10, 0.5,
         "2026-03-22", 16, LIGHT, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 슬라이드 2: 목차
# ════════════════════════════════════════
slide = make_slide()
add_text(slide, 0.6, 0.3, 6, 0.7, "목차 (Table of Contents)", 32, ACCENT, True)
add_accent_bar(slide, 0.6, 0.95, 2.5, 0.05)

sections = [
    ("01", "서론 — 분석 배경 및 목적, 데이터셋 설명"),
    ("02", "데이터 프로파일링 — 명세 확인, 결측치, 기초 통계량, 정제"),
    ("03", "단변량 분석 — 정규성 검정, 히스토그램/KDE, 클래스 분포, Binning"),
    ("04", "다변량 분석 — 상관관계, 다중공선성(VIF), 산점도, 세그먼트 비교"),
    ("05", "핵심 인사이트 — 주요 패턴, 가설 검증, 예상 밖 발견"),
    ("06", "결론 — 분석 요약, 비즈니스 제언, 한계점"),
    ("A ", "부록 — 변환 비교, 이상치 탐지, 스케일링"),
]

for idx, (num, title) in enumerate(sections):
    y = 1.4 + idx * 0.7
    # 번호
    add_text(slide, 1.0, y, 0.8, 0.5, num, 24, ACCENT, True, PP_ALIGN.RIGHT)
    # 제목
    add_text(slide, 2.0, y, 9.0, 0.5, title, 20, WHITE, False)


# ════════════════════════════════════════
# 슬라이드 3: 서론 - 분석 배경
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 1)
add_text(slide, 1.2, 0.3, 8, 0.6, "서론 — 분석 배경 및 목적", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

tf = add_text(slide, 0.8, 1.3, 12, 5.0,
              "분석 목적", 22, YELLOW, True)
add_para(tf, "• 대만 신주시 헌혈 센터의 헌혈자 데이터 기반", 18, LIGHT)
add_para(tf, "• 다음 헌혈 캠페인 시 실제로 헌혈할 가능성이 높은 헌혈자 식별", 18, WHITE, True)
add_para(tf, "• RFMTC 모델에 기반한 헌혈 행동 데이터 다각도 분석", 18, LIGHT)
add_para(tf, "• 효율적인 헌혈자 관리 및 마케팅 전략 수립에 기여", 18, LIGHT)

add_para(tf, "", 12)
add_para(tf, "주요 분석 질문 (Hypotheses)", 22, YELLOW, True, Pt(16))
add_para(tf, "H1  최근에 헌혈한 사람(V1 낮음)이 다음에도 헌혈할 확률이 높을 것이다", 18, LIGHT)
add_para(tf, "H2  헌혈 빈도(V2)가 높은 사람이 다시 헌혈할 가능성이 높을 것이다", 18, LIGHT)
add_para(tf, "H3  V3(Monetary) = V2(Frequency) × 250 완전 선형 관계일 것이다", 18, LIGHT)


# ════════════════════════════════════════
# 슬라이드 4: 데이터셋 설명
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 1)
add_text(slide, 1.2, 0.3, 8, 0.6, "서론 — 데이터셋 설명", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

# 데이터셋 개요 테이블
data_overview = [
    ["항목", "내용"],
    ["출처", "OpenML (ID: 1464), UCI ML Repository"],
    ["수집 주체", "대만 신주시 헌혈 센터"],
    ["원본 행 수", "748건"],
    ["분석 행 수", "533건 (중복 215건 제거 후)"],
    ["열 수", "5개 (특성 4 + 타겟 1)"],
]
add_table(slide, 6, 2, data_overview, 0.6, 1.2, 5.0, 2.5, font_size=13)

# 변수 설명 테이블
var_data = [
    ["변수", "의미", "설명"],
    ["V1 (Recency)", "마지막 헌혈 이후 경과 개월", "값 작을수록 최근 헌혈"],
    ["V2 (Frequency)", "총 헌혈 횟수", "헌혈 빈도"],
    ["V3 (Monetary)", "총 헌혈량 (c.c.)", "V2 × 250"],
    ["V4 (Time)", "첫 헌혈 이후 경과 개월", "헌혈 이력 기간"],
    ["Class", "2007.3 헌혈 여부", "1=미헌혈, 2=헌혈"],
]
add_table(slide, 6, 3, var_data, 6.2, 1.2, 6.5, 2.5, font_size=13)

tf = add_text(slide, 0.6, 4.5, 12, 2.0,
              "핵심 포인트", 20, YELLOW, True)
add_para(tf, "• 원본 Class 인코딩: 1(미헌혈 76.2%) / 2(헌혈 23.8%)", 16, LIGHT)
add_para(tf, "• 전체 수치형(int64), 결측치 0건 → 데이터 품질 양호", 16, LIGHT)
add_para(tf, "• 4가지 방법(openml, fetch_openml, StringIO, requests)으로 동일 결과 확인", 16, LIGHT)


# ════════════════════════════════════════
# 슬라이드 5: 결측치 & 기초 통계량
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 2)
add_text(slide, 1.2, 0.3, 8, 0.6, "데이터 프로파일링 — 기초 통계량", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

stats_data = [
    ["통계량", "V1 (Recency)", "V2 (Frequency)", "V3 (Monetary)", "V4 (Time)"],
    ["평균", "9.529", "6.991", "1747.65", "42.304"],
    ["중앙값", "8.000", "5.000", "1250.00", "38.000"],
    ["최빈값", "2", "3", "750", "28"],
    ["Q1 / Q3", "3.0 / 14.0", "3.0 / 9.0", "750 / 2250", "26.0 / 58.0"],
    ["IQR", "11.000", "6.000", "1500.00", "32.000"],
    ["분산", "68.077", "39.261", "2,453,824", "542.671"],
    ["표준편차", "8.251", "6.266", "1566.47", "23.295"],
    ["왜도", "2.265", "3.041", "3.041", "0.564"],
    ["첨도", "12.397", "13.746", "13.746", "-0.489"],
]
add_table(slide, 10, 5, stats_data, 0.5, 1.2, 12.3, 4.5, font_size=12)

tf = add_text(slide, 0.5, 6.0, 12, 1.2,
              "", 16, LIGHT)
add_para(tf, "• 평균 > 중앙값 > 최빈값 → V1~V3 양의 왜도  |  결측치 0건  |  중복 215건(28.7%) 제거 후 533건 기준", 15, LIGHT)


# ════════════════════════════════════════
# 슬라이드 6: 정규성 검정
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 3)
add_text(slide, 1.2, 0.3, 10, 0.6, "단변량 분석 — 정규성 검정 (Shapiro-Wilk)", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

norm_data = [
    ["변수", "왜도", "해석", "첨도", "해석", "Shapiro W", "p-value", "판정"],
    ["V1", "2.265", "강한 양의 왜도", "12.397", "극도 뾰족", "0.8026", "0.000", "비정규"],
    ["V2", "3.041", "강한 양의 왜도", "13.746", "극도 뾰족", "0.7236", "0.000", "비정규"],
    ["V3", "3.041", "강한 양의 왜도", "13.746", "극도 뾰족", "0.7236", "0.000", "비정규"],
    ["V4", "0.564", "중간 양의 왜도", "-0.489", "정규 유사", "0.9559", "0.000", "비정규"],
]
add_table(slide, 5, 8, norm_data, 0.4, 1.2, 12.5, 2.5, font_size=12)

tf = add_text(slide, 0.4, 3.9, 12, 3.0,
              "핵심 결론", 22, YELLOW, True)
add_para(tf, "• 4개 변수 모두 p < 0.05 → 정규분포를 따르지 않음", 18, WHITE, True)
add_para(tf, "• 비모수적 방법(Spearman 상관계수 등) 사용 필요", 18, LIGHT)
add_para(tf, "• V1~V3의 첨도 12~14 → 극도로 뾰족하여 이상치 다수 존재", 18, LIGHT)
add_para(tf, "• pandas .kurtosis()는 초과 첨도(excess kurtosis) 반환 (정규분포 기준값 = 0)", 16, RGBColor(0xA0, 0xA0, 0xA0))


# ════════════════════════════════════════
# 슬라이드 7: 히스토그램 + KDE
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 3)
add_text(slide, 1.2, 0.3, 10, 0.6, "단변량 분석 — 히스토그램 & KDE / 박스플롯", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "01_histogram_kde.png", 0.3, 1.2, width=6.3)
add_image_safe(slide, "02_boxplot.png", 6.8, 1.2, width=6.2)

tf = add_text(slide, 0.3, 6.0, 6, 1.0, "", 14, LIGHT)
add_para(tf, "녹색=평균 | 주황=중앙값 | 파랑=최빈값", 13, RGBColor(0xA0, 0xA0, 0xA0))

tf = add_text(slide, 6.8, 6.0, 6, 1.0, "", 14, LIGHT)
add_para(tf, "V2, V3: 극도의 오른쪽 꼬리 → 소수 고빈도 헌혈자", 13, RGBColor(0xA0, 0xA0, 0xA0))


# ════════════════════════════════════════
# 슬라이드 8: Class 분포 & Binning
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 3)
add_text(slide, 1.2, 0.3, 10, 0.6, "단변량 분석 — 클래스 분포 & Binning", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "03_class_distribution.png", 0.3, 1.2, width=6.0)
add_image_safe(slide, "04_binning.png", 6.5, 1.2, width=6.5)

class_data = [
    ["Class", "레이블", "빈도", "비율"],
    ["1", "미헌혈", "384", "72.0%"],
    ["2", "헌혈", "149", "28.0%"],
]
add_table(slide, 3, 4, class_data, 0.5, 5.4, 5.0, 1.2, font_size=13)

tf = add_text(slide, 6.5, 5.4, 6, 1.5, "", 16, LIGHT)
add_para(tf, "• 클래스 불균형: 72% vs 28%", 16, WHITE, True)
add_para(tf, "• 5% 기준 희소 클래스는 없음", 16, LIGHT)
add_para(tf, "• ML 시 SMOTE 등 오버샘플링 권장", 16, LIGHT)


# ════════════════════════════════════════
# 슬라이드 9: Pearson & Spearman 상관관계
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 4)
add_text(slide, 1.2, 0.3, 10, 0.6, "다변량 분석 — 상관관계 히트맵", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "07_pearson_heatmap.png", 0.2, 1.2, width=5.8, height=4.3)
add_image_safe(slide, "08_spearman_heatmap.png", 6.7, 1.2, width=5.8, height=4.3)

tf = add_text(slide, 0.2, 5.7, 12.5, 1.5, "", 15, LIGHT)
add_para(tf, "• V2-V3: Pearson/Spearman 모두 1.000 → 완벽한 다중공선성 (V3 제거 필수)", 15, WHITE, True)
add_para(tf, "• V1-Class: Spearman -0.298 → 최근성이 재헌혈의 가장 강력한 예측자  |  비정규 분포 → Spearman이 더 신뢰적", 14, LIGHT)


# ════════════════════════════════════════
# 슬라이드 10: VIF & 공분산
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 4)
add_text(slide, 1.2, 0.3, 10, 0.6, "다변량 분석 — 다중공선성 (VIF) & 공분산", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

vif_data = [
    ["변수", "VIF", "판정"],
    ["V1", "2.20", "정상"],
    ["V2", "∞ (inf)", "완벽한 다중공선성"],
    ["V3", "∞ (inf)", "완벽한 다중공선성"],
    ["V4", "5.55", "경계 수준 (기준 10 미달)"],
]
add_table(slide, 5, 3, vif_data, 0.5, 1.3, 5.0, 2.5, font_size=14)

add_image_safe(slide, "09_covariance_heatmap.png", 6.2, 1.2, width=6.5)

tf = add_text(slide, 0.5, 4.2, 5.5, 3.0,
              "다중공선성 분석 결론", 20, YELLOW, True)
add_para(tf, "• V2-V3: VIF = ∞", 18, RED, True)
add_para(tf, "  → V3 = V2 × 250 동일 정보", 16, LIGHT)
add_para(tf, "  → 모델링 시 V3 반드시 제거", 16, LIGHT)
add_para(tf, "", 10)
add_para(tf, "• V4: VIF = 5.55", 18, ORANGE, True)
add_para(tf, "  → 기준(10) 미달이나 주시 필요", 16, LIGHT)
add_para(tf, "• V1: VIF = 2.20 → 문제없음", 18, GREEN, True)


# ════════════════════════════════════════
# 슬라이드 11: 산점도 & Pair Plot
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 4)
add_text(slide, 1.2, 0.3, 10, 0.6, "다변량 분석 — 산점도 & Pair Plot", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "10c_scatter_trendline.png", 0.2, 1.2, width=5.8, height=4.3)
add_image_safe(slide, "10_pairplot.png", 6.7, 1.1, width=5.8, height=4.3)

tf = add_text(slide, 0.2, 5.7, 12, 1.5, "", 14, LIGHT)
add_para(tf, "• V2-V3: 완벽 직선 (기울기=250, R²=1.000)  |  R²<0.1인 쌍: 회색 점선(약한 관계)  |  V1-Class: 음의 추세", 14, LIGHT)


# ════════════════════════════════════════
# 슬라이드 12: 타겟별 분석
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 4)
add_text(slide, 1.2, 0.3, 10, 0.6, "다변량 분석 — Class별 분포 비교", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "11_class_comparison.png", 0.2, 1.2, width=8.0)

tf = add_text(slide, 8.5, 1.2, 4.5, 5.5,
              "Class별 주요 차이", 22, YELLOW, True)
add_para(tf, "", 8)
add_para(tf, "V1 (Recency)", 20, ACCENT, True)
add_para(tf, "Class 2(헌혈)의 중앙값이 확연히 낮음", 16, LIGHT)
add_para(tf, "→ 최근 헌혈자가 재헌혈", 16, WHITE, True)
add_para(tf, "", 8)
add_para(tf, "V2 (Frequency)", 20, ACCENT, True)
add_para(tf, "Class 2의 중앙값이 높음", 16, LIGHT)
add_para(tf, "→ 고빈도 헌혈자가 재헌혈 경향", 16, WHITE, True)
add_para(tf, "", 8)
add_para(tf, "V4 (Time)", 20, ACCENT, True)
add_para(tf, "두 Class 간 차이 미미", 16, LIGHT)
add_para(tf, "→ 이력 기간은 예측력 약함", 16, ORANGE)


# ════════════════════════════════════════
# 슬라이드 13: 세그먼트 비교
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 4)
add_text(slide, 1.2, 0.3, 10, 0.6, "다변량 분석 — VIP vs 일반 세그먼트 비교", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

add_image_safe(slide, "11b_segment_comparison.png", 0.2, 1.2, width=7.5)

seg_data = [
    ["세그먼트", "기준", "인원", "V1 평균", "V2 평균", "V4 평균"],
    ["VIP", "V2 > 5", "257명", "7.650", "11.101", "53.284"],
    ["General", "V2 ≤ 5", "276명", "11.279", "3.163", "32.080"],
]
add_table(slide, 3, 6, seg_data, 8.0, 1.4, 5.0, 1.2, font_size=12)

tf = add_text(slide, 8.0, 3.0, 5.0, 4.0,
              "시사점", 20, YELLOW, True)
add_para(tf, "• VIP: V1 평균 7.65개월", 16, GREEN, True)
add_para(tf, "  (General 11.28 대비 낮음)", 15, LIGHT)
add_para(tf, "  → 재헌혈 경향 높음", 15, LIGHT)
add_para(tf, "", 8)
add_para(tf, "• General: 비교적 최근 유입된", 16, LIGHT)
add_para(tf, "  저빈도 헌혈자 (V4 짧음)", 15, LIGHT)
add_para(tf, "", 8)
add_para(tf, "• VIP 세그먼트 집중 캠페인", 16, WHITE, True)
add_para(tf, "  → 효율 극대화 가능", 16, ACCENT, True)


# ════════════════════════════════════════
# 슬라이드 14: 핵심 인사이트
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 5)
add_text(slide, 1.2, 0.3, 10, 0.6, "핵심 인사이트 — 주요 패턴 발견", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

findings = [
    ("1", "극도의 양의 왜도", "V1~V3 모두 강한 양의 왜도(2.3~3.0)\n소수 고빈도 헌혈자가 분포 왜곡", ORANGE),
    ("2", "V2-V3 완벽 상관", "Pearson/Spearman 모두 1.000, VIF = ∞\n동일 정보 → V3 제거 필수", RED),
    ("3", "V1이 가장 강력한 예측자", "Class와 Spearman -0.298\n비정규 분포 고려 시 Spearman 기준 신뢰", GREEN),
    ("4", "클래스 불균형", "헌혈 28.0% vs 미헌혈 72.0%\nML 시 오버샘플링 필요", YELLOW),
    ("5", "V4의 낮은 예측력", "Class와 Spearman -0.142\n오랜 이력 ≠ 재헌혈 보장", LIGHT),
]

for idx, (num, title, desc, color) in enumerate(findings):
    x = 0.5 + (idx % 3) * 4.2
    y = 1.3 + (idx // 3) * 2.8

    # 카드 배경
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(3.8), Inches(2.4))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    add_text(slide, x + 0.2, y + 0.15, 3.4, 0.4, f"#{num}", 16, color, True)
    add_text(slide, x + 0.2, y + 0.5, 3.4, 0.4, title, 18, WHITE, True)

    lines = desc.split('\n')
    for li, line in enumerate(lines):
        add_text(slide, x + 0.2, y + 1.0 + li * 0.45, 3.4, 0.4, line, 14, LIGHT)


# ════════════════════════════════════════
# 슬라이드 15: 가설 검증
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 5)
add_text(slide, 1.2, 0.3, 10, 0.6, "핵심 인사이트 — 가설 검증 결과", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

hyp_data = [
    ["가설", "내용", "결과", "근거"],
    ["H1", "최근 헌혈자가 재헌혈", "채택 ✓", "V1-Class Spearman -0.298, Box Plot 확인"],
    ["H2", "고빈도 헌혈자가 재헌혈", "채택 ✓", "V2-Class Pearson 0.175, VIP V1 평균 7.65 < General 11.28"],
    ["H3", "V3 = V2 × 250", "채택 ✓", "Pearson/Spearman 1.000, 산점도 완벽 직선, VIF = ∞"],
]
add_table(slide, 4, 4, hyp_data, 0.5, 1.4, 12.3, 2.2, font_size=14)

tf = add_text(slide, 0.5, 4.0, 12, 3.0,
              "예상치 못한 발견 (Unexpected Findings)", 22, YELLOW, True)
add_para(tf, "", 6)
add_para(tf, "1. 중복 데이터 28.7%  — 748건 중 215건이 완전 동일 RFMTC 패턴", 18, LIGHT)
add_para(tf, "   → 동일 행동 양식의 헌혈자가 다수 존재", 16, RGBColor(0xA0, 0xA0, 0xA0))
add_para(tf, "", 6)
add_para(tf, "2. Modified Z-score vs Z-score 큰 괴리  — V2에서 Z 9개 vs MAD 39개", 18, LIGHT)
add_para(tf, "   → 비정규 분포에서 Z-score가 이상치 과소 탐지", 16, RGBColor(0xA0, 0xA0, 0xA0))
add_para(tf, "", 6)
add_para(tf, "3. VIP의 역설  — VIP(V2>5)의 V4 평균 53.3개월 > General 32.1개월", 18, LIGHT)
add_para(tf, "   → 오래 헌혈 = 고빈도 (자명), 그러나 V4 자체는 Class 예측에 약함", 16, RGBColor(0xA0, 0xA0, 0xA0))


# ════════════════════════════════════════
# 슬라이드 16: 결론
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 6)
add_text(slide, 1.2, 0.3, 10, 0.6, "결론 — 분석 요약 & 비즈니스 제언", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

# 왼쪽: 분석 요약
tf = add_text(slide, 0.5, 1.3, 6.0, 6.0,
              "분석 요약", 22, ACCENT, True)
add_para(tf, "", 6)
add_para(tf, "• 533건 대상, 35개 기술통계기법 EDA 수행", 17, LIGHT)
add_para(tf, "• 결측치 0건, 양호한 데이터 품질", 17, LIGHT)
add_para(tf, "• 극도의 양의 왜도 (V1~V3)", 17, LIGHT)
add_para(tf, "• 클래스 불균형 (72:28)", 17, LIGHT)
add_para(tf, "• V2-V3 완전 다중공선성 (VIF = ∞)", 17, LIGHT)
add_para(tf, "• V1(Recency)이 가장 유의미한 예측 변수", 17, WHITE, True)
add_para(tf, "  (Spearman -0.298)", 15, RGBColor(0xA0, 0xA0, 0xA0))

# 오른쪽: 제언
tf = add_text(slide, 7.0, 1.3, 6.0, 6.0,
              "비즈니스 제언", 22, ACCENT, True)
add_para(tf, "", 6)
add_para(tf, "1. 타겟 마케팅", 18, YELLOW, True)
add_para(tf, "   V1↓ V2↑ VIP 세그먼트(257명) 집중 캠페인", 16, LIGHT)
add_para(tf, "", 6)
add_para(tf, "2. 변수 선택", 18, YELLOW, True)
add_para(tf, "   V3 제거 → V1, V2, V4 3개 변수 모델링", 16, LIGHT)
add_para(tf, "", 6)
add_para(tf, "3. 변환 적용", 18, YELLOW, True)
add_para(tf, "   박스-콕스: V1 왜도 2.265→-0.013", 16, LIGHT)
add_para(tf, "", 6)
add_para(tf, "4. 불균형 처리", 18, YELLOW, True)
add_para(tf, "   SMOTE 오버샘플링으로 Class 균형", 16, LIGHT)
add_para(tf, "", 6)
add_para(tf, "5. 모델 추천", 18, YELLOW, True)
add_para(tf, "   트리 기반 (Random Forest, XGBoost)", 16, LIGHT)


# ════════════════════════════════════════
# 슬라이드 17: 한계점
# ════════════════════════════════════════
slide = make_slide()
add_section_number(slide, 6)
add_text(slide, 1.2, 0.3, 10, 0.6, "결론 — 한계점 및 추후 과제", 30, WHITE, True)
add_accent_bar(slide, 1.2, 0.85, 2.0)

limits = [
    ("변수 부족", "4개(실질 3개) 변수만으로는\n예측 정확도 한계", "📊"),
    ("시계열 미반영", "헌혈 패턴의 시간적 추이\n분석 부재", "📈"),
    ("외부 변수 부재", "연령, 성별, 건강 상태 등\n미포함", "👤"),
    ("이상치 처리 미수행", "탐지만 하고 제거/변환\n하지 않음 (EDA 목적상 원본 유지)", "⚠"),
]

for idx, (title, desc, icon) in enumerate(limits):
    x = 0.6 + idx * 3.15
    y = 1.5

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y),
                                   Inches(2.9), Inches(3.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x40)
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.5)

    add_text(slide, x + 0.3, y + 0.3, 2.3, 0.5, title, 18, ORANGE, True)
    lines = desc.split('\n')
    for li, line in enumerate(lines):
        add_text(slide, x + 0.3, y + 0.9 + li * 0.45, 2.3, 0.4, line, 15, LIGHT)

tf = add_text(slide, 0.6, 5.0, 12, 2.0,
              "향후 방향", 22, GREEN, True)
add_para(tf, "추가 변수 수집  +  앙상블 모델링  +  머신러닝 도입으로 예측 고도화", 18, WHITE, True)


# ════════════════════════════════════════
# 슬라이드 18: 부록 - 변환 비교
# ════════════════════════════════════════
slide = make_slide()
add_text(slide, 0.4, 0.3, 4, 0.6, "부록 A", 18, ACCENT, True)
add_text(slide, 1.6, 0.3, 10, 0.6, "— 왜곡 변수 전처리 (변환 비교)", 30, WHITE, True)
add_accent_bar(slide, 1.6, 0.85, 2.5)

add_image_safe(slide, "05_transformations.png", 0.2, 1.2, width=6.5)

trans_data = [
    ["변수", "원본 왜도", "로그", "루트", "박스-콕스", "Yeo-Johnson", "최적"],
    ["V1", "2.265", "-0.076", "0.506", "-0.013", "-0.013", "박스-콕스"],
    ["V2", "3.041", "0.405", "1.301", "0.002", "0.020", "박스-콕스"],
    ["V3", "3.041", "0.084", "1.301", "0.002", "0.002", "박스-콕스"],
    ["V4", "0.564", "-0.972", "-0.059", "-0.050", "-0.048", "Yeo-Johnson"],
]
add_table(slide, 5, 7, trans_data, 7.0, 1.3, 6.0, 2.5, font_size=11)

add_image_safe(slide, "05b_yeojohnson.png", 7.0, 4.2, width=6.0)


# ════════════════════════════════════════
# 슬라이드 19: 부록 - 이상치 탐지
# ════════════════════════════════════════
slide = make_slide()
add_text(slide, 0.4, 0.3, 4, 0.6, "부록 A", 18, ACCENT, True)
add_text(slide, 1.6, 0.3, 10, 0.6, "— 이상치 탐지 & 스케일링", 30, WHITE, True)
add_accent_bar(slide, 1.6, 0.85, 2.5)

add_image_safe(slide, "06_outliers.png", 0.2, 1.2, width=6.5)

outlier_data = [
    ["변수", "Z-score (|Z|>3)", "Modified Z (|M|>3.5)", "IQR 방식"],
    ["V1", "6개", "3개", "6개"],
    ["V2", "9개", "39개", "21개"],
    ["V3", "9개", "39개", "21개"],
    ["V4", "0개", "0개", "0개"],
]
add_table(slide, 5, 4, outlier_data, 7.0, 1.3, 5.5, 2.5, font_size=12)

add_image_safe(slide, "12_scaling_comparison.png", 7.0, 4.2, width=5.5)

tf = add_text(slide, 7.0, 3.9, 5.5, 0.5,
              "비정규 분포 → MAD 기반 Modified Z-score가 더 강건", 14, ORANGE)


# ════════════════════════════════════════
# 슬라이드 20: 부록 - CDF & 프로파일
# ════════════════════════════════════════
slide = make_slide()
add_text(slide, 0.4, 0.3, 4, 0.6, "부록 B", 18, ACCENT, True)
add_text(slide, 1.6, 0.3, 10, 0.6, "— 추가 시각화 (CDF & Class 프로파일)", 30, WHITE, True)
add_accent_bar(slide, 1.6, 0.85, 2.5)

add_image_safe(slide, "13_cdf_lineplot.png", 0.2, 1.2, width=6.3)
add_image_safe(slide, "14_class_profile_lineplot.png", 6.8, 1.2, width=6.2)

add_image_safe(slide, "10b_violin_overall.png", 2.5, 5.0, width=8.0)


# ════════════════════════════════════════
# 슬라이드 21: Thank You
# ════════════════════════════════════════
slide = make_slide()
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0),
                                SLIDE_W, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, 1.5, 2.2, 10, 1.0,
         "감사합니다", 48, WHITE, True, PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(5.0), Inches(3.4),
                                Inches(3.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_text(slide, 1.5, 3.8, 10, 0.6,
         "Q & A", 32, ACCENT, True, PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.8, 10, 0.5,
         "혈액 기부 서비스 센터 데이터 탐색적 분석 (EDA)", 18, LIGHT, False, PP_ALIGN.CENTER)


# ── 저장 ──
prs.save(OUT)
print(f"PPT 생성 완료: {OUT}")
print(f"총 슬라이드 수: {len(prs.slides)}")
